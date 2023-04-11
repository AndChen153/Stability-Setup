import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches
import numpy as np
from pptx.enum.text import PP_ALIGN
import os

# Create a new PowerPoint presentation
class PowerPointCreator:
    '''
    Class to create powerpoints for collected data
    '''

    def __init__(self, PnOcsv:str, scanCSV: str, PnOimage: str, scanLightBefore: str,scanLightAfter: str,scanDarkBefore: str,scanDarkAfter: str, experimentTitle: str, leadName: str) -> None:
        self.prs = Presentation()
        # self.prs.slide_width = Inches(16)
        # self.prs.slide_height = Inches(9)
        # dir_list = os.listdir(path)
        self.PnOcsv = PnOcsv
        self.pno = []
        self.scanLightBefore = []
        self.scanLightAfter = []
        self.scanDarkBefore = []
        self.scanDarkAfter = []

        for i in os.listdir(PnOimage):
            self.pno.append(PnOimage+"\\"+i)

        for i in os.listdir(scanLightBefore):
            self.scanLightBefore.append(scanLightBefore+"\\"+i)

        for i in os.listdir(scanLightAfter):
            self.scanLightAfter.append(scanLightAfter+"\\"+i)

        for i in os.listdir(scanDarkBefore):
            self.scanDarkBefore.append(scanDarkBefore+"\\"+i)

        for i in os.listdir(scanDarkAfter):
            self.scanDarkAfter.append(scanDarkAfter+"\\"+i)




        self.experiement = experimentTitle

        # create title slide
        title_slide_layout = self.prs.slide_layouts[0]
        slide              = self.prs.slides.add_slide(title_slide_layout)
        title              = slide.shapes.title
        subtitle           = slide.placeholders[1]

        title.text = self.experiement
        # title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle.text = leadName

        self.slide1Title()
        self.slide3PNO()
        # self.slide3_1()
        self.slide4Metrics()
        self.slide5DarkJV()
        self.slide6LightJV()


    def slide1Title(self):
        slide_layout = self.prs.slide_layouts[1]
        measurementParams = np.loadtxt(self.PnOcsv, delimiter=",", dtype=str)[0:5,0:2]
        # startDate = self.pno["csv"].split('\\')[-1].split(" ")[0][3:]

        slide = self.prs.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title

        title_shape.text = 'PnO     Measurement Conditions'

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = np.array2string(measurementParams).translate({ord(i): None for i in '[]\'\''})


        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = np.array2string(measurementParams).translate({ord(i): None for i in '[]\'\''})
        # tf.text += "\n Start Date:"
        # tf.text += startDate


    def slide2Metrics(self):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = 'Scan  Metrics'
        device0PNO = slide.shapes.add_picture(self.pno[0], Inches(0.69), Inches(1.38), Inches(4.31)) #horizontal pos, vert pos, width
        device1PNO = slide.shapes.add_picture(self.pno[1], Inches(5), Inches(1.38), Inches(4.31)) #horizontal pos, vert dispost, width
        device2PNO = slide.shapes.add_picture(self.pno[2], Inches(0.69), Inches(4.41), Inches(4.31)) #horizontal pos, vert pos, width
        device3PNO = slide.shapes.add_picture(self.pno[3], Inches(5), Inches(4.41), Inches(4.31)) #horizontal pos, vert pos, width


    def slide3PNO(self):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = 'Averaged PNO Traces (no dead cells)'
        device0PNO = slide.shapes.add_picture(self.pno[0], Inches(0.69), Inches(1.38), Inches(4.31)) #horizontal pos, vert pos, width
        device1PNO = slide.shapes.add_picture(self.pno[1], Inches(5), Inches(1.38), Inches(4.31)) #horizontal pos, vert dispost, width
        device2PNO = slide.shapes.add_picture(self.pno[2], Inches(0.69), Inches(4.41), Inches(4.31)) #horizontal pos, vert pos, width
        device3PNO = slide.shapes.add_picture(self.pno[3], Inches(5), Inches(4.41), Inches(4.31)) #horizontal pos, vert pos, width

    def slide4Metrics(self):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = 'Scan  params'
        voc = slide.shapes.add_picture(self.scanDarkBefore[-1], Inches(3.46), Inches(2.64), Inches(3.08)) #horizontal pos, vert pos, width
        jsc = slide.shapes.add_picture(self.scanDarkBefore[-2], Inches(0.24), Inches(2.64), Inches(3.08)) #horizontal pos, vert dispost, width
        ff = slide.shapes.add_picture(self.scanDarkBefore[-3], Inches(6.68), Inches(2.64), Inches(3.08)) #horizontal pos, vert pos, width




    def slide5DarkJV(self):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = 'Dark JV Before |  Dark JV After'

        device0BeforeDark = slide.shapes.add_picture(self.scanDarkBefore[0], Inches(0), Inches(1.92), Inches(2.5)) #left dist, top dist, width
        device1BeforeDark = slide.shapes.add_picture(self.scanDarkBefore[1], Inches(2.5), Inches(1.92), Inches(2.5)) #left dist, top dist, height
        device2BeforeDark = slide.shapes.add_picture(self.scanDarkBefore[2], Inches(0), Inches(3.82), Inches(2.5)) #left dist, top dist, height
        device3BeforeDark = slide.shapes.add_picture(self.scanDarkBefore[3], Inches(2.5), Inches(3.82), Inches(2.5)) #left dist, top dist, height

        device0AfterDark = slide.shapes.add_picture(self.scanDarkAfter[0], Inches(5), Inches(1.92), Inches(2.5)) #left dist, top dist, height
        device1AfterDark = slide.shapes.add_picture(self.scanDarkAfter[1], Inches(7.5), Inches(1.92), Inches(2.5)) #left dist, top dist, height
        device2AfterDark = slide.shapes.add_picture(self.scanDarkAfter[2], Inches(5), Inches(3.82), Inches(2.5)) #left dist, top dist, height
        device3AfterDark = slide.shapes.add_picture(self.scanDarkAfter[3], Inches(7.5), Inches(3.82), Inches(2.5)) #left dist, top dist, height



    def slide6LightJV(self):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = 'Light JV Before |   Light JV After'

        device0BeforeLight = slide.shapes.add_picture(self.scanLightBefore[0], Inches(0), Inches(1.92), Inches(2.5)) #left dist, top dist, width
        device1BeforeLight = slide.shapes.add_picture(self.scanLightBefore[1], Inches(2.5), Inches(1.92), Inches(2.5)) #left dist, top dist, height
        device2BeforeLight = slide.shapes.add_picture(self.scanLightBefore[2], Inches(0), Inches(3.82), Inches(2.5)) #left dist, top dist, height
        device3BeforeLight = slide.shapes.add_picture(self.scanLightBefore[3], Inches(2.5), Inches(3.82), Inches(2.5)) #left dist, top dist, height

        device0AfterLight = slide.shapes.add_picture(self.scanLightAfter[0], Inches(5), Inches(1.92), Inches(2.5)) #left dist, top dist, height
        device1AfterLight = slide.shapes.add_picture(self.scanLightAfter[1], Inches(7.5), Inches(1.92), Inches(2.5)) #left dist, top dist, height
        device2AfterLight = slide.shapes.add_picture(self.scanLightAfter[2], Inches(5), Inches(3.82), Inches(2.5)) #left dist, top dist, height
        device3AfterLight = slide.shapes.add_picture(self.scanLightAfter[3], Inches(7.5), Inches(3.82), Inches(2.5)) #left dist, top dist, height


    def save(self) -> str:
        loc = "./dataVisualization/" + self.experiement + ".pptx"
        self.prs.save(loc)
        print("saved at " + loc)
        return loc


if __name__ == '__main__':
    pno = {"imagePnO":r"C:\Users\achen\Dropbox\code\Stability-Setup\data\Mar-15-2023\PnOMar-15-2023 13_29_55\\",
           "csv": r"C:\Users\achen\Dropbox\code\Stability-Setup\data\Mar-15-2023\PnOMar-15-2023 13_29_55.csv"}
    scan = {"imageLight":r"C:\Users\achen\Dropbox\code\Stability-Setup\data\Mar-15-2023\scanlightMar-15-2023 13_28_50\\",
            "csvLight":r"C:\Users\achen\Dropbox\code\Stability-Setup\data\Mar-15-2023\scanlightMar-15-2023 13_28_50.csv",
            "imageDark":r"C:\Users\achen\Dropbox\code\Stability-Setup\data\Mar-15-2023\scandarkMar-15-2023 13_34_51\\",
            "csvDark":r"C:\Users\achen\Dropbox\code\Stability-Setup\data\Mar-15-2023\scandarkMar-15-2023 13_34_51.csv"}
    ppt = PowerPointCreator(pno["csv"], scan["csvLight"],
                            pno["imagePnO"],
                            scan["imageLight"],scan["imageLight"],
                            scan["imageDark"],scan["imageDark"],
                            "Example Auto PPT", "Andrew Chen")
    ppt.save()
